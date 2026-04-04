from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt


OUTPUT = Path("outputs/ROOKHIDE_Hackathon_Ready_Pitch.pptx")

FONT_NAME = "Calibri"
TITLE_SIZE = Pt(38)
SUBTITLE_SIZE = Pt(24)
BODY_SIZE = Pt(22)
SMALL_SIZE = Pt(16)
ACCENT = RGBColor(187, 67, 45)
TEXT_COLOR = RGBColor(25, 25, 25)


SLIDES = [
    {
        "title": "Hide Any File Inside a Valid Chess Game",
        "subtitle": "Private file sharing that looks like ordinary gameplay",
        "points": [
            "ROOKHIDE converts file data into legal chess move choices",
            "Output remains valid PGN and is replayable in chess tools",
        ],
        "notes": "Open with one line: private files hidden in normal-looking chess games.",
    },
    {
        "title": "Real-World Problem",
        "subtitle": "Sensitive documents are shared daily on common channels",
        "points": [
            "Offer letters, IDs, contracts, and invoices travel via chat, email, and cloud links",
            "Encrypted files can raise suspicion in practical non-technical contexts",
            "Existing privacy tools often add friction and poor usability",
            "Need: privacy that blends into normal digital behavior",
        ],
        "notes": "Keep it relatable; no jargon.",
    },
    {
        "title": "Solution",
        "subtitle": "Use legal chess branching as the data channel",
        "points": [
            "Each legal chess position offers multiple valid moves",
            "ROOKHIDE maps bits to deterministic move selection",
            "Compression and optional password protection run before encoding",
            "Result: standard PGN that looks like a real game",
        ],
        "notes": "Explain the core insight clearly with a visual board snapshot.",
    },
    {
        "title": "How It Works",
        "subtitle": "End-to-end pipeline",
        "points": [
            "Upload file -> compress/encrypt payload",
            "Convert payload to bits and encode via legal moves",
            "Generate one or more valid chess PGN games",
            "Decode by replaying move order to recover original bytes",
        ],
        "notes": "Show a left-to-right process diagram.",
    },
    {
        "title": "Feasibility",
        "subtitle": "Technically and operationally practical",
        "points": [
            "Tech: Flask orchestration + Rust engine via PyO3",
            "Operational: browser-based interface, simple deployment",
            "Security: optional encryption, deterministic recovery",
            "Performance: around 1 MB processed in about 3.6 seconds",
        ],
        "notes": "Mention clear module separation: UI, backend, engine.",
    },
    {
        "title": "Business Model (INR)",
        "subtitle": "Subscription plus enterprise licensing",
        "points": [
            "Plans: Free, Pro Rs 299/month, Team Rs 1,499/month, Enterprise Rs 14,999/month",
            "Estimated monthly cost: Hosting Rs 15,000 + Support Rs 20,000 + Marketing Rs 10,000",
            "Estimated operating cost: about Rs 45,000/month",
            "Sample MRR: 100 Pro + 20 Team + 5 Enterprise = Rs 1,34,875/month",
        ],
        "notes": "Keep pricing and unit economics visible in one chart.",
    },
    {
        "title": "Market Opportunity",
        "subtitle": "Target users and proven spend in privacy tech",
        "points": [
            "Target users: students, freelancers, HR teams, law firms, fintech teams, enterprises",
            "India cybersecurity market: roughly Rs 25,000-30,000 crore annually",
            "Global encryption/privacy-tech revenue: roughly Rs 1.2-1.3 lakh crore annually",
            "Existing encryption products show strong recurring willingness to pay",
        ],
        "notes": "Present TAM-SAM-SOM quickly with audience segments.",
    },
    {
        "title": "Penetration Testing and Resilience",
        "subtitle": "Existing approaches vs ROOKHIDE under detection and tamper scenarios",
        "points": [
            "AES/RSA file artifacts are typically flagged by entropy and signature checks",
            "Image/audio steganography is often exposed by statistical steganalysis and ML models",
            "ROOKHIDE outputs valid PGN, reducing file-format detectability in baseline scans",
            "Internal red-team checks: detection attempts, tamper tests, replay consistency, decode integrity",
        ],
        "notes": "Use a comparison matrix and clarify that claims are based on internal tests and baseline scanners.",
    },
    {
        "title": "Team (4 Members)",
        "subtitle": "Clear ownership and contribution",
        "points": [
            "Anurag: Product direction and backend orchestration",
            "Kartik: Testing, documentation, and deployment support",
            "Pranay: Frontend, visualizer, and user experience",
            "Rishu: Rust engine and performance optimization",
        ],
        "notes": "Use team photo and four role cards.",
    },
    {
        "title": "ROOKHIDE",
        "subtitle": "Hide anything. Look normal.",
        "points": [
            "Relatable problem",
            "Detailed, feasible solution",
            "Clear business and market case",
            "Hackathon-ready team and demo",
        ],
        "notes": "Close with QR code and demo link.",
    },
]


def style_paragraph(paragraph, size, bold=False, color=TEXT_COLOR):
    for run in paragraph.runs:
        run.font.name = FONT_NAME
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    title_p = slide.shapes.title.text_frame.paragraphs[0]
    style_paragraph(title_p, TITLE_SIZE, bold=True)

    sub_p = slide.placeholders[1].text_frame.paragraphs[0]
    style_paragraph(sub_p, SUBTITLE_SIZE)

    return slide


def add_content_slide(prs, title, subtitle, points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    title_p = slide.shapes.title.text_frame.paragraphs[0]
    style_paragraph(title_p, Pt(34), bold=True)

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    first = tf.paragraphs[0]
    first.text = subtitle
    style_paragraph(first, SUBTITLE_SIZE, bold=True, color=ACCENT)

    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        style_paragraph(p, BODY_SIZE)

    hint = slide.shapes.add_textbox(Pt(40), Pt(485), Pt(840), Pt(36))
    hint_tf = hint.text_frame
    hint_tf.clear()
    hint_p = hint_tf.paragraphs[0]
    hint_p.text = "Visual hint: keep this slide image-heavy (1 chart/screenshot + short bullets)."
    style_paragraph(hint_p, SMALL_SIZE, color=RGBColor(110, 110, 110))

    return slide


def add_notes(slide, notes_text):
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes_text


def build_presentation():
    prs = Presentation()

    for idx, data in enumerate(SLIDES):
        if idx == 0:
            slide = add_title_slide(prs, data["title"], data["subtitle"])
        else:
            slide = add_content_slide(prs, data["title"], data["subtitle"], data["points"])
        add_notes(slide, data.get("notes", ""))

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
    print(f"Created: {OUTPUT}")
